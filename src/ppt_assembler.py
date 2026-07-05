"""
ppt_assembler.py
Places chart images into NIIF template placeholder slots.
Preserves all template branding — logo, header, colors, everything.
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from src.chart_builder import build_chart
from src.ppt_logger    import get_logger, log_exceptions

R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

CHART_PH_TYPES = {
    PP_PLACEHOLDER.CHART,    # 8
    PP_PLACEHOLDER.OBJECT,   # 7
    PP_PLACEHOLDER.PICTURE,  # 9
    PP_PLACEHOLDER.TABLE,    # 12 — some single-chart layouts use TABLE slot
}


def _get_chart_slots(layout):
    """Get chart placeholder positions sorted top→bottom, left→right.
    Falls back to largest BODY placeholder if no chart/table slots found."""
    slots = []
    for ph in layout.placeholders:
        pf = ph.placeholder_format
        if pf.type in CHART_PH_TYPES:
            slots.append({
                "idx":    pf.idx,
                "left":   ph.left,
                "top":    ph.top,
                "width":  ph.width,
                "height": ph.height,
            })
    if not slots:
        body_phs = []
        for ph in layout.placeholders:
            pf = ph.placeholder_format
            if pf.type == PP_PLACEHOLDER.BODY:
                body_phs.append((ph.width * ph.height, ph, pf))
        if body_phs:
            body_phs.sort(key=lambda x: -x[0])
            ph, pf = body_phs[0][1], body_phs[0][2]
            slots.append({
                "idx":    pf.idx,
                "left":   ph.left,
                "top":    ph.top,
                "width":  ph.width,
                "height": ph.height,
            })
    slots.sort(key=lambda s: (s["top"], s["left"]))
    return slots


def _fallback_slots(n, sw_emu, sh_emu):
    """Fallback positions when template has no chart placeholders."""
    M  = int(0.25 * 914400)
    T  = int(1.20 * 914400)
    G  = int(0.15 * 914400)
    B  = int(0.35 * 914400)
    uw = sw_emu - M * 2
    uh = sh_emu - T - B
    if n == 1:
        return [{"left": M, "top": T, "width": uw, "height": uh}]
    elif n == 2:
        cw = (uw - G) // 2
        return [
            {"left": M,       "top": T, "width": cw, "height": uh},
            {"left": M+cw+G,  "top": T, "width": cw, "height": uh},
        ]
    elif n == 3:
        cw = (uw - G*2) // 3
        return [{"left": M+i*(cw+G), "top": T, "width": cw, "height": uh}
                for i in range(3)]
    else:
        cw = (uw - G) // 2
        ch = (uh - G) // 2
        return [
            {"left": M,      "top": T,      "width": cw, "height": ch},
            {"left": M+cw+G, "top": T,      "width": cw, "height": ch},
            {"left": M,      "top": T+ch+G, "width": cw, "height": ch},
            {"left": M+cw+G, "top": T+ch+G, "width": cw, "height": ch},
        ]


def _remove_all_slides(prs):
    count = len(prs.slides._sldIdLst)
    for _ in range(count):
        el  = prs.slides._sldIdLst[0]
        rId = el.get(f"{{{R_NS}}}id")
        if rId:
            try: prs.part.drop_rel(rId)
            except: pass
        del prs.slides._sldIdLst[0]
    return count


def _get_layout(prs, name):
    """Search all slide masters for layout by name — exact then partial match."""
    name_up     = name.upper().strip()
    all_layouts = []
    for master in prs.slide_masters:
        all_layouts.extend(master.slide_layouts)
    all_layouts.extend(prs.slide_layouts)

    for l in all_layouts:
        if l.name.upper().strip() == name_up:
            return l
    for l in all_layouts:
        if name_up in l.name.upper():
            return l

    print(f"  WARNING: layout '{name}' not found — using first available")
    print(f"  Available: {[l.name for l in all_layouts[:10]]}")
    return all_layouts[0] if all_layouts else prs.slide_layouts[0]


def _set_title(slide, text):
    """Fill title placeholder if present."""
    for ph in slide.placeholders:
        if ph.placeholder_format.type in (
            PP_PLACEHOLDER.TITLE,
            PP_PLACEHOLDER.CENTER_TITLE,
        ):
            try:
                ph.text = text
                tf = ph.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    run = tf.paragraphs[0].runs[0]
                    run.font.size = Pt(18)
                    run.font.bold = True
            except Exception:
                pass
            break


def _set_subtitle(slide, text):
    """Fill the subtitle BODY placeholder (default text 'Click to edit Master text styles')."""
    if not text:
        return
    layout = slide.slide_layout
    # find the subtitle placeholder idx from layout
    subtitle_idx = None
    for ph in layout.placeholders:
        pf = ph.placeholder_format
        if pf.type == PP_PLACEHOLDER.BODY:
            try:
                default_txt = ph.text_frame.text.strip()
            except Exception:
                continue
            if default_txt.startswith("Click to edit Master text"):
                subtitle_idx = pf.idx
                break
    if subtitle_idx is None:
        return
    # fill on the actual slide
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == subtitle_idx:
            try:
                ph.text = text
                tf = ph.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    run = tf.paragraphs[0].runs[0]
                    run.font.size = Pt(12)
            except Exception:
                pass
            break


def _fill_headers(slide, slide_cfg, positions):
    """Fill Chart/Table Header, Sub-Header, and Source placeholders
    by matching layout default text + (left, top) rank ordering."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER

    layout = slide.slide_layout

    # ── Step 1: Group layout BODY placeholders by their default text ──
    groups = {"header": [], "subheader": [], "source": [],
              "note": [], "content": []}

    for ph in layout.placeholders:
        pf = ph.placeholder_format
        if pf.type != PP_PLACEHOLDER.BODY:
            continue
        try:
            txt = ph.text_frame.text.strip()
        except Exception:
            continue

        first_line = txt.split("\n")[0].strip()
        entry = {
            "idx":  pf.idx,
            "left": ph.left / 914400,
            "top":  ph.top  / 914400,
        }

        if first_line == "Chart/Table Header":
            groups["header"].append(entry)
        elif first_line.startswith("Chart/Table Sub"):
            groups["subheader"].append(entry)
        elif first_line == "Source":
            groups["source"].append(entry)
        elif first_line.startswith("Note"):
            groups["note"].append(entry)
        elif first_line in ("Content", "Comment"):
            groups["content"].append(entry)

    # ── Step 2: Sort each group by (left, top) — round to avoid float noise ──
    for g in groups.values():
        g.sort(key=lambda x: (round(x["left"], 1), round(x["top"], 1)))

    # ── Step 3: Sort chart positions by (left, top), keep original index ──
    indexed_positions = sorted(
        enumerate(positions),
        key=lambda x: (round(x[1]["left"] / 914400, 1),
                        round(x[1]["top"]  / 914400, 1)),
    )

    # ── Step 4: Build slide placeholder lookup by idx ──
    slide_ph_map = {}
    for ph in slide.placeholders:
        slide_ph_map[ph.placeholder_format.idx] = ph

    # Debug: show what was found
    print(f"    [GROUPS] header={len(groups['header'])} subheader={len(groups['subheader'])} "
          f"source={len(groups['source'])} note={len(groups['note'])} "
          f"content={len(groups['content'])}")
    for gname, glist in groups.items():
        if glist:
            idxs = [g["idx"] for g in glist]
            on_slide = [i for i in idxs if i in slide_ph_map]
            missing  = [i for i in idxs if i not in slide_ph_map]
            if missing:
                print(f"    [WARN] {gname}: layout has idx={idxs} "
                      f"but slide missing idx={missing}")

    # Force-create missing placeholders from layout
    from pptx.oxml.ns import qn
    from copy import deepcopy
    layout_ph_map = {}
    for ph in layout.placeholders:
        layout_ph_map[ph.placeholder_format.idx] = ph
    for gname, glist in groups.items():
        for g in glist:
            idx = g["idx"]
            if idx not in slide_ph_map and idx in layout_ph_map:
                # clone the layout placeholder XML onto the slide
                layout_el = layout_ph_map[idx]._element
                new_el = deepcopy(layout_el)
                slide.shapes._spTree.append(new_el)
                # re-read into slide_ph_map
                for ph in slide.placeholders:
                    if ph.placeholder_format.idx == idx:
                        slide_ph_map[idx] = ph
                        break

    # ── Step 5: Match rank-by-rank and fill ──
    print(f"    [HEADERS] {len(groups['header'])} header slots, "
          f"{len(indexed_positions)} chart positions")
    for rank, (chart_i, pos) in enumerate(indexed_positions):
        if chart_i >= len(slide_cfg.charts):
            continue
        chart_cfg = slide_cfg.charts[chart_i]

        # ── Header ──
        # ── Header ──
        header = getattr(chart_cfg, "chart_header", "").strip()
        print(f"    [HEADER] rank={rank} chart[{chart_i}] "
              f"pos=({pos['left']/914400:.1f}\",{pos['top']/914400:.1f}\") "
              f"→ '{header[:50]}'")
        if header and rank < len(groups["header"]):
            hdr_entry = groups["header"][rank]
            print(f"             → placeholder idx={hdr_entry['idx']} "
                  f"at ({hdr_entry['left']:.1f}\",{hdr_entry['top']:.1f}\")")
            ph = slide_ph_map.get(groups["header"][rank]["idx"])
            if ph:
                try:
                    tf = ph.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text      = header
                    run.font.name = "Calibri"
                    run.font.size = Pt(14)
                    run.font.bold = True
                except Exception:
                    pass

        # ── Sub-header ──
        subheader = getattr(chart_cfg, "chart_subheader", "").strip()
        if subheader and rank < len(groups["subheader"]):
            ph = slide_ph_map.get(groups["subheader"][rank]["idx"])
            if ph:
                try:
                    tf = ph.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text      = subheader
                    run.font.name = "Calibri"
                    run.font.size = Pt(14)
                    run.font.bold = True
                except Exception:
                    pass

        # ── Source ──
        source = getattr(chart_cfg, "chart_source", "").strip()
        if not source:
            source = getattr(chart_cfg, "chart_source_ai", "").strip()
        if source and rank < len(groups["source"]):
            ph = slide_ph_map.get(groups["source"][rank]["idx"])
            if ph:
                try:
                    tf = ph.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text           = f"Source: {source}"
                    run.font.size      = Pt(7)
                    run.font.italic    = True
                    run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
                except Exception:
                    pass

        # ── Note ──
        note = getattr(chart_cfg, "chart_note", "").strip()
        if note and rank < len(groups["note"]):
            ph = slide_ph_map.get(groups["note"][rank]["idx"])
            if ph:
                try:
                    tf = ph.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text           = f"Note: {note}"
                    run.font.size      = Pt(7)
                    run.font.italic    = True
                    run.font.color.rgb = RGBColor(0x7F, 0x7F, 0x7F)
                except Exception:
                    pass

        # ── Content / Comment ──
        content_val = getattr(chart_cfg, "chart_content", "").strip()
        if content_val:
            print(f"    [CONTENT] rank={rank} chart[{chart_i}] "
                  f"val='{content_val[:40]}' "
                  f"content_slots={len(groups['content'])}")
            if rank < len(groups["content"]):
                cidx = groups["content"][rank]["idx"]
                print(f"             → trying placeholder idx={cidx} "
                      f"on_slide={cidx in slide_ph_map}")
        if content_val and rank < len(groups["content"]):
            ph = slide_ph_map.get(groups["content"][rank]["idx"])
            if ph:
                try:
                    tf = ph.text_frame
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text      = content_val
                    run.font.size = Pt(7)
                except Exception:
                    pass

@log_exceptions
def assemble_ppt(slides, template_path, output_path, chart_theme):
    prs    = Presentation(template_path)
    sw_emu = prs.slide_width
    sh_emu = prs.slide_height

    removed = _remove_all_slides(prs)
    print(f"  Removed {removed} template slide(s). "
          f"Building {len(slides)} new slide(s)...")

    _log = get_logger()
    for slide_cfg in slides:
        _log.set_context(slide=slide_cfg.slide_number, chart="")
        layout = _get_layout(prs, slide_cfg.layout_name)
        slide  = prs.slides.add_slide(layout)
        slots  = _get_chart_slots(layout)
        n      = slide_cfg.chart_count

        # ── Slide start banner ───────────────────────────────────────
        _heading    = getattr(slide_cfg, "slide_heading",     "") or slide_cfg.slide_title or ""
        _subheading = getattr(slide_cfg, "slide_sub_heading", "") or ""
        print(f"\n{'='*60}")
        print(f"  SLIDE {slide_cfg.slide_number}  |  {n} chart(s)  |  layout: {layout.name}")
        print(f"  Heading    : {_heading or '(none)'}")
        print(f"  Sub-heading: {_subheading or '(none)'}")
        print(f"{'='*60}")

        # slide title: prefer AI slide_heading, fall back to config title
        title_text = getattr(slide_cfg, "slide_heading", "") or slide_cfg.slide_title
        _set_title(slide, title_text)
        # slide subtitle
        _set_subtitle(slide, getattr(slide_cfg, "slide_sub_heading", ""))

        # pick slot positions
        if len(slots) >= n:
            positions = slots[:n]
            print(f"    Using template placeholder positions")
        else:
            positions = _fallback_slots(min(n, 4), sw_emu, sh_emu)
            print(f"    Using calculated positions "
                  f"(template has {len(slots)} slots, need {n})")

        # remove empty chart placeholders (replaced by actual images below)
        chart_ph_idxs = {s["idx"] for s in positions if "idx" in s}
        sp_tree       = slide.shapes._spTree
        to_remove     = [
            ph._element for ph in slide.placeholders
            if ph.placeholder_format.type in CHART_PH_TYPES
            and ph.placeholder_format.idx in chart_ph_idxs
        ]
        for el in to_remove:
            try: sp_tree.remove(el)
            except: pass

        # render and place each chart
        for i, chart_cfg in enumerate(slide_cfg.charts[:len(positions)]):
            _log.set_context(chart=chr(65 + i))
            pos = positions[i]
            fw  = round(pos["width"]  / 914400, 2)
            fh  = round(pos["height"] / 914400, 2)
            try:
                buf = build_chart(chart_cfg, fig_width=fw, fig_height=fh,theme=chart_theme)
                buf.seek(0)
                slide.shapes.add_picture(
                    buf,
                    pos["left"], pos["top"],
                    pos["width"], pos["height"])
                print(f"    Chart {i+1}: OK  "
                      f"pos=({round(pos['left']/914400,1)}\", "
                      f"{round(pos['top']/914400,1)}\")  "
                      f"{fw:.1f}\" x {fh:.1f}\"")
            except Exception as e:
                print(f"    Chart {i+1}: ERROR — {e}")
                _log = get_logger()
                _log.exception(
                    f"Chart render failed — slide {slide_cfg.slide_number} chart {i+1}",
                    chart_type=getattr(chart_cfg, "y1_chart_type", "?"),
                    widget_id=getattr(chart_cfg, "widget_id", "?"),
                    fw=fw, fh=fh,
                )

        # fill header and source placeholders
        _fill_headers(slide, slide_cfg, positions)
        print(f"  ✔  Slide {slide_cfg.slide_number} done")

    prs.save(output_path)
    print(f"\n✅  Saved → {output_path}")
    return output_path